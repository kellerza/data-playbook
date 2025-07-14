"""Powerpoint helpers."""

from __future__ import annotations
import re
import typing as t
from typing import Any

import attrs
from colordict import ColorDict
from icecream import ic

# from pptx import Presentation as NewPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.text.text import TextFrame, _Paragraph
from pptx.util import Length, Pt

RE_STYLES = re.compile(r"(.*?)(?:<([A-Z,0-9#-]+)>|$)")


def int_length(val: float | Length | None) -> Length | None:
    """Convert int to Length."""
    if isinstance(val, Length):
        return val
    if isinstance(val, int | float):
        return Pt(val)
    if val is None:
        return None

    return Length(val)


@attrs.define
class PStyle:
    """Paragraph style class."""

    bold: bool | None = None
    color: RGBColor | tuple[int, int, int] | None = None
    highlight: RGBColor | tuple[int, int, int] | None = None
    italic: bool | None = None
    size: Length | None = attrs.field(default=None, converter=int_length)
    strike: bool | None = None

    def __bool__(self) -> bool:
        """Check if the style has any attributes set."""
        return bool(
            self.bold
            or self.color
            or self.highlight
            or self.italic
            or self.size
            or self.strike
        )


class PText:
    """Paragraph text class."""

    _list: list

    def __init__(self, *text: str | PStyle | PText):
        """Init Paragraph text."""
        self._list = []
        if text:
            self.append(*text)

    def append(self, *text: str | PStyle | PText) -> None:
        """Parse values and append."""
        for val in text:
            if isinstance(val, PStyle):
                self._list.append(val)
                continue
            if isinstance(val, PText):
                self.append(*val._list)
                continue
            if not isinstance(val, str):
                raise ValueError(f"expected string, got {type(val)}")

            for line in val.splitlines(keepends=True):
                # _LOGGER.info("'%s'", line)
                newline = line.endswith("\n")
                matches = RE_STYLES.findall(line.strip("\n"))
                for run in matches:
                    if run[0]:
                        self._list.append(run[0])
                    if run[1]:
                        self._list.append(str2styles(run[1]))
                if newline:
                    self._list.append("\n")

    def apply_to(self, para: _Paragraph) -> None:
        """Apply to, was SlidePlaceholders."""
        style: PStyle | None = None
        for prun in self._list:
            if isinstance(prun, str):
                if prun == "\n":
                    para.add_line_break()
                    continue
                run = para.add_run()
                # https://python-pptx.readthedocs.io/en/latest/api/text.html#run-objects
                run.text = prun
                if style:
                    if style.highlight:
                        rpr = run._r.get_or_add_rPr()
                        hl = OxmlElement("a:highlight")
                        clr = OxmlElement("a:srgbClr")
                        clr.val = str(style.highlight)  # type:ignore[attr-defined]
                        hl.append(clr)
                        rpr.append(hl)  # type:ignore[arg-type]
                    if style.size:
                        run.font.size = style.size
                    if style.color:
                        run.font.color.rgb = style.color
                    run.font.italic = style.italic
                    run.font.bold = style.bold
                    if style.strike:
                        run.font.underline = MSO_TEXT_UNDERLINE_TYPE.DOT_DASH_HEAVY_LINE
                        # https://github.com/scanny/python-pptx/pull/606/files
                        # rPr = run.get_or_add_rPr()
                        rpr = run.font._rPr
                        rpr.strike = "sngStrike"  # type:ignore[attr-defined]
                        # rPr.strike = "sngStrike"
                        # rPr.strikethrough = "sngStrike"
                    # run.hyperlink
                style = None
            if isinstance(prun, PStyle):
                style = prun


def str2styles(style_s: str) -> PStyle:
    """Styles."""
    style_s = style_s.strip().strip("!")
    res = PStyle()
    if not style_s:
        return res
    ss = style_s.split(",")

    res.bold = True if "B" in ss else None
    res.italic = True if "I" in ss else None
    res.strike = True if "-" in ss else None
    for rem in ("B", "I", "-"):
        while rem in ss:
            ss.remove(rem)

    for kk in list(ss):
        try:
            # ic("size", kk, float(kk), Pt(float(kk)))
            res.size = Pt(int(kk))
            ss.remove(kk)
        except (TypeError, ValueError):
            continue

    highlight = False
    # _LOGGER.error("conv col '%s'", ss)
    for col in ss:
        if col == "":
            highlight = True
            continue

        col = col.lower()  # noqa: PLW2901
        # Convert string to color
        if not col.startswith("#"):
            try:
                col = str(ColorDict(mode="hex")[col])  # noqa: PLW2901
            except Exception as err:
                ic("BAD color", col, err)
        try:
            rgb = RGBColor.from_string(col.strip("#"))
            if highlight:
                res.highlight = rgb
            else:
                res.color = rgb
        except Exception as err:
            ic("BAD hex color", col, err)

        highlight = True

    return res


def add_paragraphs(
    *,
    frame: TextFrame,
    text: Any,
    clear: bool = True,
    indent: int = 0,
) -> None:
    """Add paragraph."""
    if not isinstance(frame, TextFrame):
        ic(".text_frame???")
        frame = frame.text_frame
    if clear:
        frame.clear()
    for ptext in text:
        if isinstance(ptext, list | tuple):
            add_paragraphs(
                frame=frame,
                text=ptext,
                clear=False,
                indent=indent + 1,
            )
            continue
        if isinstance(ptext, str):
            ptext = PText(ptext)  # noqa: PLW2901
        if not isinstance(ptext, PText):
            ic("string/PText expected", ptext)
            ptext = PText(str(ptext))  # noqa: PLW2901

        # get the paragraph - if clear, paragraph[0] exist, but it's empty
        para = frame.paragraphs[0] if clear else frame.add_paragraph()
        clear = False

        if indent > 0:
            para.level = indent

        ptext.apply_to(para)


def add_slide(prs: Presentation, layout: str) -> Slide:
    """Add & return a roadmap slide."""
    slide_layouts = [lay.name for lay in prs.slide_layouts]
    try:
        lay = prs.slide_layouts[slide_layouts.index(layout)]
        # lay = prs.slides[TEMPLATE_ROADMAP_NO].slide_layout
    except Exception as err:
        slide_layouts = [lay.name for lay in prs.slide_layouts]
        ic("available layouts / templates", slide_layouts)
        raise err
    return prs.slides.add_slide(lay)


@t.runtime_checkable
class ShapeWithText(t.Protocol):
    """Shape with text."""

    has_text_frame: bool
    text: str
    top: int


@attrs.define
class Slide3Parts:
    """Analyse a slide & the 3 main parts."""

    slide: Slide
    shapes: dict[int, ShapeWithText] = attrs.field(init=False, default={})

    def __attrs_post_init__(self) -> None:
        """Init the class."""
        self.get_shapes()

    def get_shapes(self) -> None:
        """Get the shape parts in order."""
        shapes: list[ShapeWithText] = []
        for sh in self.slide.shapes:
            if isinstance(sh, ShapeWithText) and sh.has_text_frame:
                shapes.append(sh)

        self.shapes = dict(enumerate(sorted(shapes, key=lambda s: s.top)))

    @property
    def title(self) -> ShapeWithText:
        """Slide title shape."""
        return self.shapes[0]

    @property
    def subtitle(self) -> ShapeWithText:
        """Slide subtitle shape."""
        return self.shapes[1]

    @property
    def text(self) -> ShapeWithText:
        """Slide main text element."""
        return self.shapes[2]

    @property
    def notes(self) -> TextFrame | None:
        """Comments."""
        return self.slide.notes_slide.notes_text_frame

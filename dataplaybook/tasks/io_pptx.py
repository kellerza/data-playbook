"""Powerpoint helpers."""
# pylint: disable=invalid-name
from __future__ import annotations

import re
from dataclasses import dataclass
from turtle import Shape
from typing import Any, Optional, Union

from colordict import ColorDict
from icecream import ic
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.slide import Slide, SlidePlaceholders
from pptx.text.text import TextFrame
from pptx.util import Pt

RE_STYLES = re.compile(r"(.*?)(?:<([A-Z,0-9#-]+)>|$)")


@dataclass
class PStyle:
    """Paragraph style class."""

    bold: Optional[bool] = None
    strike: Optional[bool] = None
    italic: Optional[bool] = None
    color: Optional[RGBColor] = None
    highlight: Optional[RGBColor] = None
    size: Optional[Pt] = None

    # def __bool__(self):
    #     return self.bold or self.italic or self.color or self.size or False


class PText:
    """Paragraph text class."""

    _list: list

    def __init__(self, *text: Union[str, PStyle]):
        """Init Paragraph text."""
        self._list = []
        if text:
            self.append(*text)

    def append(self, *text: Union[str, PStyle]) -> None:
        """Parse values and append."""
        for val in text:
            if isinstance(val, PStyle):
                self._list.append(val)
                continue
            if isinstance(val, PText):
                self.append(*val._list)  # pylint:disable=protected-access
                continue
            if not isinstance(val, str):
                raise ValueError(f"expected string, got {type(val)}")

            for line in val.splitlines(keepends=True):
                # _LOGGER.info("'%s'", line)
                newline = line.endswith("\n")
                matches = RE_STYLES.findall(line.strip("\n"))
                for run in matches:
                    if run[0]:
                        self._list.append(run[0])
                    if run[1]:
                        self._list.append(str2styles(run[1]))
                if newline:
                    self._list.append("\n")

    def apply_to(self, para: SlidePlaceholders) -> None:
        """Apply to."""
        style: Optional[PStyle] = None
        for prun in self._list:
            if isinstance(prun, str):
                if prun == "\n":
                    para.add_line_break()
                    continue
                run = para.add_run()
                # https://python-pptx.readthedocs.io/en/latest/api/text.html#run-objects
                run.text = prun
                if style:
                    if style.highlight:
                        rPr = run._r.get_or_add_rPr()  # pylint:disable=protected-access
                        hl = OxmlElement("a:highlight")
                        srgbClr = OxmlElement("a:srgbClr")
                        setattr(srgbClr, "val", str(style.highlight))
                        hl.append(srgbClr)
                        rPr.append(hl)
                    if style.size:
                        run.font.size = Pt(style.size)
                    if style.color:
                        run.font.color.rgb = style.color
                    run.font.italic = style.italic
                    run.font.bold = style.bold
                    if style.strike:
                        run.font.underline = MSO_TEXT_UNDERLINE_TYPE.DOT_DASH_HEAVY_LINE
                        # https://github.com/scanny/python-pptx/pull/606/files
                        # rPr = run.get_or_add_rPr()
                        rPr = run.font._rPr  # pylint:disable=protected-access
                        rPr.strike = "sngStrike"
                        # rPr.strikethrough = "sngStrike"
                    # run.hyperlink
                style = None
            if isinstance(prun, PStyle):
                style = prun


def str2styles(style_s: str) -> PStyle:
    """Styles."""
    style_s = style_s.strip().strip("!")
    res = PStyle()
    if not style_s:
        return res
    ss = style_s.split(",")

    res.bold = True if "B" in ss else None
    res.italic = True if "I" in ss else None
    res.strike = True if "-" in ss else None
    for rem in ("B", "I", "-"):
        while rem in ss:
            ss.remove(rem)

    for kk in list(ss):
        try:
            res.size = int(kk)
            ss.remove(kk)
        except (TypeError, ValueError):
            continue

    highlight = False
    # _LOGGER.error("conv col '%s'", ss)
    for col in ss:
        if col == "":
            highlight = True
            continue

        col = col.lower()
        # Convert string to color
        if not col.startswith("#"):
            try:
                col = ColorDict(mode="hex")[col]
            except Exception as err:  # noqa pylint: disable=broad-exception-caught
                ic("BAD color", col, err)
        try:
            rgb = RGBColor.from_string(col.strip("#"))
            if highlight:
                res.highlight = rgb
            else:
                res.color = rgb
        except Exception as err:  # pylint: disable=broad-exception-caught
            ic("BAD hex color", col, err)

        highlight = True

    return res


def add_paragraphs(
    *,
    frame: TextFrame,
    text: Any,
    clear: bool = True,
    indent: int = 0,
) -> None:
    """Add paragraph."""
    if not isinstance(frame, TextFrame):
        ic(".text_frame???")
        frame = frame.text_frame
    if clear:
        frame.clear()
    for ptext in text:
        if isinstance(ptext, (list, tuple)):
            add_paragraphs(
                frame=frame,
                text=ptext,
                clear=False,
                indent=indent + 1,
            )
            continue
        if isinstance(ptext, str):
            ptext = PText(ptext)
        if not isinstance(ptext, PText):
            ic("string/PText expected", ptext)
            ptext = PText(str(ptext))

        # get the paragraph - if clear, paragraph[0] exist, but it's empty
        para = frame.paragraphs[0] if clear else frame.add_paragraph()
        clear = False

        if indent > 0:
            para.level = indent

        ptext.apply_to(para)


def add_slide(prs: Presentation, layout: str) -> Slide:
    """Add & return a roadmap slide."""
    slide_layouts = [lay.name for lay in prs.slide_layouts]
    try:
        lay = prs.slide_layouts[slide_layouts.index(layout)]
        # lay = prs.slides[TEMPLATE_ROADMAP_NO].slide_layout
    except Exception as err:
        slide_layouts = [lay.name for lay in prs.slide_layouts]
        ic("available layouts / templates", slide_layouts)
        raise err
    return prs.slides.add_slide(lay)


class Slide3Parts:
    """Analyse a slide & the 3 main parts."""

    _sl: Slide = None
    _sh: dict[int, Shape] = {}

    def __init__(self, sl: Slide):
        """Init the class."""
        self._sl = sl
        self._sh: dict[int, Shape] = {}
        self.get_sh()

    def get_sh(self) -> None:
        """Get the shape parts in order."""
        sh: SlidePlaceholders
        for sh in self._sl.shapes:
            if not sh.has_text_frame:
                continue
            self._sh[sh.top] = sh

        for idx, top in enumerate(sorted(self._sh.keys())):
            self._sh[idx] = self._sh.pop(top)

    @property
    def title(self) -> Shape:
        """Slide title shape."""
        return self._sh[0]

    @property
    def subtitle(self) -> Shape:
        """Slide subtitle shape."""
        return self._sh[1]

    @property
    def text(self) -> Shape:
        """Slide main text element."""
        return self._sh[2]

    @property
    def notes(self) -> TextFrame:
        """Comments."""
        return self._sl.notes_slide.notes_text_frame
